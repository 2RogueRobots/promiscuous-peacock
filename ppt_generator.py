"""
PowerPoint Generator Module for Promiscuous-Peacock

A reusable module for creating professional PowerPoint presentations from
notebook content with smart layout handling based on image aspect ratios.

Template: test.pptx (APONTIS branded)
- Slide size: 13.33" x 7.5" (16:9 widescreen)
- Font: Reddit Sans
- Brand colors: #FF325D (primary), #2D185C (text), #FCFCFC (contrast)

Usage:
    from ppt_generator import PPTGenerator, PPTConfig

    config = PPTConfig(footer_text="APONTIS | Confidential")
    ppt = PPTGenerator("test.pptx", config)
    ppt.add_title_slide("Market Analysis", "Q1 2026 Report")
    ppt.add_image_slide("REGIONAL DATA", "NRW leads market", "chart.png")
    ppt.save("output.pptx")
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Literal, TYPE_CHECKING

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

if TYPE_CHECKING:
    import pandas as pd

# Try to import PIL for image detection, fall back gracefully
try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False


# =============================================================================
# BRAND CONSTANTS
# =============================================================================

BRAND = {
    "font": "Reddit Sans",
    "primary": "#FF325D",      # Red - slide titles on content slides
    "text_dark": "#2D185C",    # Dark purple - body text on white bg
    "text_light": "#FCFCFC",   # White - text on dark bg (title/end slides)
}

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_LEFT = Inches(0.63)
MARGIN_RIGHT = Inches(0.63)
CONTENT_WIDTH = Inches(11.99)

# =============================================================================
# ELEMENT SPECIFICATIONS
# =============================================================================

TITLE_SPEC = {
    "left": Inches(0.63),
    "top": Inches(0.2),
    "width": Inches(11.99),
    "height": Inches(0.41),
    "font_size": Pt(28),
    "caps": True,
}

ACTION_TITLE_SPEC = {
    "left": Inches(0.63),
    "top": Inches(0.61),
    "width": Inches(11.99),
    "height": Inches(0.47),
    "font_size": Pt(20),
    "caps": False,
}

BODY_SPEC = {
    "left": Inches(0.63),
    "top": Inches(1.2),
    "width": Inches(11.99),
    "height": Inches(4.5),
    "font_size": Pt(14),
    "line_spacing": 1.15,
}

FOOTER_SPEC = {
    "left": Inches(2.82),
    "top": Inches(6.95),
    "width": Inches(8.21),
    "height": Inches(0.4),
    "font_size": Pt(10),
}

# =============================================================================
# IMAGE LAYOUT POSITIONS
# =============================================================================

IMAGE_LAYOUTS = {
    "image_top": {
        "image": {"left": 0.63, "top": 1.2, "width": 11.99, "height": 4.0},
        "text": {"left": 0.63, "top": 5.4, "width": 11.99, "height": 1.3},
    },
    "image_left": {
        "image": {"left": 0.63, "top": 1.2, "width": 5.5, "height": 4.5},
        "text": {"left": 6.5, "top": 1.2, "width": 5.5, "height": 4.5},
    },
    "image_right": {
        "image": {"left": 6.5, "top": 1.2, "width": 5.5, "height": 4.5},
        "text": {"left": 0.63, "top": 1.2, "width": 5.5, "height": 4.5},
    },
    "full_image": {
        "image": {"left": 0.5, "top": 1.2, "width": 12.33, "height": 5.3},
    },
}

# Template layout indices
LAYOUT_SECTION_HEADER = 0
LAYOUT_TITLE_SLIDE = 1
LAYOUT_BLANK = 2
LAYOUT_TITLE_VARIANT = 3
LAYOUT_TITLE_AND_CONTENT = 4


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RgbColor."""
    hex_color = hex_color.lstrip("#")
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def detect_image_layout(image_path: str | Path) -> str:
    """
    Determine optimal layout based on image aspect ratio.

    Returns:
        'image_top': Landscape (wide) images - ratio > 1.5
        'image_right': Portrait (tall) images - ratio < 0.75
        'image_left': Square-ish images - 0.75 <= ratio <= 1.5
    """
    if not HAS_PIL:
        # Default to image_top if PIL not available
        return "image_top"

    path = Path(image_path)
    if not path.exists():
        return "image_top"

    try:
        with Image.open(path) as img:
            ratio = img.width / img.height
    except Exception:
        return "image_top"

    if ratio > 1.5:
        return "image_top"
    elif ratio < 0.75:
        return "image_right"
    else:
        return "image_left"


# =============================================================================
# CONFIGURATION
# =============================================================================

@dataclass
class PPTConfig:
    """
    Global configuration for presentations.

    Args:
        footer_text: Required footer text for all slides
        title_font: Font for titles (default: Reddit Sans)
        body_font: Font for body text (default: Reddit Sans)
        title_caps: Whether to uppercase all titles (default: True)
        text_color: Hex color for text (default: #2D185C)
        accent_color: Hex color for accents (default: #FF325D)
    """
    footer_text: str
    title_font: str = BRAND["font"]
    body_font: str = BRAND["font"]
    title_caps: bool = True
    text_color: str = BRAND["text_dark"]
    accent_color: str = BRAND["primary"]
    text_light: str = BRAND["text_light"]

    def __post_init__(self):
        if not self.footer_text:
            raise ValueError("footer_text is required and cannot be empty")


# =============================================================================
# SLIDE BUILDER (FLUENT API)
# =============================================================================

class SlideBuilder:
    """
    Fluent API for building slide content.

    Example:
        builder = ppt.add_content_slide("MARKET ANALYSIS")
        builder.action_title("Key findings").body("Details here").image("chart.png")
    """

    def __init__(self, slide, config: PPTConfig, generator: "PPTGenerator"):
        self._slide = slide
        self._config = config
        self._generator = generator
        self._has_image = False
        self._image_layout = None

    def title(self, text: str, caps: bool | None = None) -> "SlideBuilder":
        """Set the main title (usually already set)."""
        if caps is None:
            caps = self._config.title_caps

        display_text = text.upper() if caps else text

        txBox = self._slide.shapes.add_textbox(
            TITLE_SPEC["left"],
            TITLE_SPEC["top"],
            TITLE_SPEC["width"],
            TITLE_SPEC["height"],
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = display_text
        p.font.name = self._config.title_font
        p.font.size = TITLE_SPEC["font_size"]
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(self._config.text_color)

        return self

    def action_title(self, text: str) -> "SlideBuilder":
        """Add action title below main title."""
        txBox = self._slide.shapes.add_textbox(
            ACTION_TITLE_SPEC["left"],
            ACTION_TITLE_SPEC["top"],
            ACTION_TITLE_SPEC["width"],
            ACTION_TITLE_SPEC["height"],
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self._config.title_font
        p.font.size = ACTION_TITLE_SPEC["font_size"]
        p.font.color.rgb = hex_to_rgb(self._config.text_color)

        return self

    def body(self, text: str, position: dict | None = None) -> "SlideBuilder":
        """
        Add body text to the slide.

        If an image has been added, text position is adjusted based on layout.
        """
        if position is None:
            if self._image_layout and "text" in IMAGE_LAYOUTS.get(self._image_layout, {}):
                pos = IMAGE_LAYOUTS[self._image_layout]["text"]
                position = {k: Inches(v) for k, v in pos.items()}
            else:
                position = {
                    "left": BODY_SPEC["left"],
                    "top": BODY_SPEC["top"],
                    "width": BODY_SPEC["width"],
                    "height": BODY_SPEC["height"],
                }

        txBox = self._slide.shapes.add_textbox(
            position["left"],
            position["top"],
            position["width"],
            position["height"],
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        # Handle multiline text
        lines = text.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.name = self._config.body_font
            p.font.size = BODY_SPEC["font_size"]
            p.font.color.rgb = hex_to_rgb(self._config.text_color)

        return self

    def image(
        self,
        path: str | Path,
        layout: Literal["auto", "image_top", "image_left", "image_right", "full_image"] = "auto",
    ) -> "SlideBuilder":
        """
        Add an image to the slide with smart positioning.

        Args:
            path: Path to the image file
            layout: Layout mode - 'auto' detects from aspect ratio
        """
        path = Path(path)
        if not path.exists():
            raise FileNotFoundError(f"Image not found: {path}")

        if layout == "auto":
            layout = detect_image_layout(path)

        self._image_layout = layout
        self._has_image = True

        pos = IMAGE_LAYOUTS[layout]["image"]

        self._slide.shapes.add_picture(
            str(path),
            Inches(pos["left"]),
            Inches(pos["top"]),
            width=Inches(pos["width"]),
        )

        return self

    def table(
        self,
        data: "pd.DataFrame",
        columns: list[str] | None = None,
        max_rows: int = 15,
    ) -> "SlideBuilder":
        """
        Add a data table to the slide.

        Args:
            data: DataFrame to display
            columns: Columns to include (default: all)
            max_rows: Maximum rows to display (default: 15)
        """
        import pandas as pd

        if columns:
            data = data[columns]

        if len(data) > max_rows:
            data = data.head(max_rows)

        rows = len(data) + 1  # +1 for header
        cols = len(data.columns)

        # Calculate table position (below action title)
        table_left = BODY_SPEC["left"]
        table_top = BODY_SPEC["top"]
        table_width = BODY_SPEC["width"]
        table_height = Inches(min(4.5, 0.3 * rows))

        shape = self._slide.shapes.add_table(
            rows, cols,
            table_left, table_top,
            table_width, table_height,
        )
        table = shape.table

        # Style headers
        for col_idx, col_name in enumerate(data.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            p = cell.text_frame.paragraphs[0]
            p.font.name = self._config.body_font
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(self._config.text_color)

        # Fill data
        for row_idx, (_, row) in enumerate(data.iterrows(), start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)

                # Format numbers
                if isinstance(value, (int, float)):
                    if pd.isna(value):
                        cell.text = ""
                    elif isinstance(value, float) and value == int(value):
                        cell.text = f"{int(value):,}"
                    elif isinstance(value, float):
                        cell.text = f"{value:,.1f}"
                    else:
                        cell.text = f"{value:,}"
                else:
                    cell.text = str(value) if not pd.isna(value) else ""

                p = cell.text_frame.paragraphs[0]
                p.font.name = self._config.body_font
                p.font.size = Pt(10)
                p.font.color.rgb = hex_to_rgb(self._config.text_color)

        return self

    def footer(self, text: str | None = None) -> "SlideBuilder":
        """Add footer to the slide (uses config footer if not specified)."""
        footer_text = text or self._config.footer_text

        txBox = self._slide.shapes.add_textbox(
            FOOTER_SPEC["left"],
            FOOTER_SPEC["top"],
            FOOTER_SPEC["width"],
            FOOTER_SPEC["height"],
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = footer_text
        p.font.name = self._config.body_font
        p.font.size = FOOTER_SPEC["font_size"]
        p.font.color.rgb = hex_to_rgb(self._config.text_color)
        p.alignment = PP_ALIGN.CENTER

        return self

    @property
    def slide(self):
        """Return the underlying slide object for advanced customization."""
        return self._slide


# =============================================================================
# MAIN GENERATOR CLASS
# =============================================================================

class PPTGenerator:
    """
    Main PowerPoint generator class.

    Creates professional presentations from template with brand styling.

    Example:
        ppt = PPTGenerator("test.pptx", config)
        ppt.add_title_slide("Analysis Report", "Q1 2026")
        ppt.add_content_slide("KEY FINDINGS").body("Finding 1\nFinding 2")
        ppt.save("output.pptx")
    """

    def __init__(self, template_path: str | Path, config: PPTConfig, clear_template: bool = True):
        """
        Initialize generator with template and configuration.

        Args:
            template_path: Path to the PowerPoint template
            config: PPTConfig with footer_text (required) and optional settings
            clear_template: Remove existing slides from template (default: True)
        """
        self._template_path = Path(template_path)
        self._config = config

        if not self._template_path.exists():
            raise FileNotFoundError(f"Template not found: {self._template_path}")

        self._prs = Presentation(str(self._template_path))

        # Ensure correct slide dimensions
        self._prs.slide_width = SLIDE_WIDTH
        self._prs.slide_height = SLIDE_HEIGHT

        # Remove existing template slides to start fresh
        if clear_template:
            self._clear_existing_slides()

    def _clear_existing_slides(self):
        """Remove all existing slides from the template."""
        # Access the XML structure to remove slides properly
        # We need to remove slide relationship IDs from the presentation part
        xml_slides = self._prs.slides._sldIdLst
        slides_to_remove = list(xml_slides)
        for sldId in slides_to_remove:
            xml_slides.remove(sldId)

    def _get_blank_layout(self):
        """Get the blank layout from template."""
        return self._prs.slide_layouts[LAYOUT_BLANK]

    def _add_title_to_slide(self, slide, title: str, color: str | None = None) -> None:
        """Add formatted title to any slide."""
        display_text = title.upper() if self._config.title_caps else title
        title_color = color or self._config.text_color

        txBox = slide.shapes.add_textbox(
            TITLE_SPEC["left"],
            TITLE_SPEC["top"],
            TITLE_SPEC["width"],
            TITLE_SPEC["height"],
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = display_text
        p.font.name = self._config.title_font
        p.font.size = TITLE_SPEC["font_size"]
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(title_color)

    def _add_footer_to_slide(self, slide, color: str | None = None) -> None:
        """Add footer to slide."""
        footer_color = color or self._config.text_color
        txBox = slide.shapes.add_textbox(
            FOOTER_SPEC["left"],
            FOOTER_SPEC["top"],
            FOOTER_SPEC["width"],
            FOOTER_SPEC["height"],
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = self._config.footer_text
        p.font.name = self._config.body_font
        p.font.size = FOOTER_SPEC["font_size"]
        p.font.color.rgb = hex_to_rgb(footer_color)
        p.alignment = PP_ALIGN.CENTER

    def add_title_slide(self, title: str, subtitle: str = "") -> SlideBuilder:
        """
        Add a title slide (dark purple background with white text).

        IMPORTANT: Use this ONCE at the very beginning of the presentation.
        This is the opening "bookend" slide.

        Args:
            title: Main title text
            subtitle: Subtitle text (optional)

        Returns:
            SlideBuilder for further customization
        """
        slide = self._prs.slides.add_slide(self._prs.slide_layouts[LAYOUT_TITLE_SLIDE])
        light_color = self._config.text_light

        # Title - WHITE text on dark background
        if slide.shapes.title:
            slide.shapes.title.text = title.upper() if self._config.title_caps else title
            for p in slide.shapes.title.text_frame.paragraphs:
                p.font.name = self._config.title_font
                p.font.color.rgb = hex_to_rgb(light_color)

        # Subtitle - WHITE text on dark background
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Subtitle placeholder
                shape.text = subtitle
                for p in shape.text_frame.paragraphs:
                    p.font.name = self._config.body_font
                    p.font.color.rgb = hex_to_rgb(light_color)
                break

        self._add_footer_to_slide(slide, color=light_color)

        return SlideBuilder(slide, self._config, self)

    def add_content_slide(self, title: str, action_title: str = "") -> SlideBuilder:
        """
        Add a content slide (white background with red title).

        IMPORTANT: Use this for ALL content. All presentation content
        goes on white background slides.

        Uses RED title and dark purple body text on white background.

        Args:
            title: Main title (will be UPPERCASED by default)
            action_title: Secondary title line (optional)

        Returns:
            SlideBuilder for adding content (body, image, table)
        """
        slide = self._prs.slides.add_slide(self._get_blank_layout())

        # RED title on content slides
        self._add_title_to_slide(slide, title, color=self._config.accent_color)

        builder = SlideBuilder(slide, self._config, self)

        if action_title:
            builder.action_title(action_title)

        # Dark purple footer on content slides
        self._add_footer_to_slide(slide, color=self._config.text_color)

        return builder

    def add_logo_slide(self) -> SlideBuilder:
        """
        Add the logo/end slide (dark purple background with template logo).

        IMPORTANT: Use this ONCE at the very end of the presentation.
        This is the closing "bookend" slide. It shows the template's logo
        only - no custom text is added.

        Returns:
            SlideBuilder for further customization (though typically none needed)
        """
        slide = self._prs.slides.add_slide(self._prs.slide_layouts[LAYOUT_SECTION_HEADER])
        # No text added - template logo shows through
        # No footer needed on end slide
        return SlideBuilder(slide, self._config, self)

    def add_image_slide(
        self,
        title: str,
        image_path: str | Path,
        action_title: str = "",
        text: str = "",
        layout: Literal["auto", "image_top", "image_left", "image_right", "full_image"] = "auto",
    ) -> SlideBuilder:
        """
        Add a slide with image and optional text.

        Layout is auto-detected based on image aspect ratio:
        - Landscape (ratio > 1.5): Image on top, text below
        - Portrait (ratio < 0.75): Image on right, text on left
        - Square: Image on left, text on right

        Args:
            title: Slide title (UPPERCASED by default)
            image_path: Path to image file
            action_title: Secondary title (optional)
            text: Body text (optional)
            layout: Layout mode - 'auto' for automatic detection

        Returns:
            SlideBuilder for further customization
        """
        builder = self.add_content_slide(title, action_title)
        builder.image(image_path, layout)

        if text:
            builder.body(text)

        return builder

    def add_table_slide(
        self,
        title: str,
        data: "pd.DataFrame",
        action_title: str = "",
        columns: list[str] | None = None,
        max_rows: int = 15,
    ) -> SlideBuilder:
        """
        Add a slide with a data table.

        Args:
            title: Slide title
            data: DataFrame to display
            action_title: Secondary title (optional)
            columns: Columns to include (default: all)
            max_rows: Maximum rows to display (default: 15)

        Returns:
            SlideBuilder for further customization
        """
        builder = self.add_content_slide(title, action_title)
        builder.table(data, columns, max_rows)

        return builder

    def save(self, output_path: str | Path) -> Path:
        """
        Save the presentation to file.

        Args:
            output_path: Output file path

        Returns:
            Path to saved file
        """
        output_path = Path(output_path)
        self._prs.save(str(output_path))
        return output_path

    @property
    def presentation(self):
        """Return underlying Presentation object for advanced use."""
        return self._prs


# =============================================================================
# CONVENIENCE FUNCTION
# =============================================================================

def quick_ppt(
    template: str | Path,
    slides: list[dict],
    output: str | Path,
    footer: str,
    **config_kwargs,
) -> Path:
    """
    Quick one-liner PowerPoint generation.

    Slide Structure (Bookend Pattern):
    - Title slide: Dark purple BG, used ONCE at beginning
    - Content slides: White BG, ALL content goes here
    - Logo slide: Dark purple BG, used ONCE at end (no custom text)

    Args:
        template: Path to template file
        slides: List of slide definitions, each dict contains:
            - type: 'title', 'content', 'image', 'table', 'logo'
            - title: Slide title (not used for 'logo' type)
            - subtitle: For title slides
            - action_title: Secondary title
            - text/body: Body text
            - image: Image path (for image slides)
            - data: DataFrame (for table slides)
            - columns: Column list (for table slides)
        output: Output file path
        footer: Footer text (required)
        **config_kwargs: Additional PPTConfig options

    Returns:
        Path to saved file

    Example:
        quick_ppt(
            template='test.pptx',
            slides=[
                {'type': 'title', 'title': 'Report', 'subtitle': 'Q1 2026'},
                {'type': 'content', 'title': 'KEY FINDINGS', 'body': 'Details...'},
                {'type': 'logo'},  # End slide - no text needed
            ],
            output='report.pptx',
            footer='Report Title Here'
        )
    """
    config = PPTConfig(footer_text=footer, **config_kwargs)
    ppt = PPTGenerator(template, config)

    for slide_def in slides:
        slide_type = slide_def.get("type", "content")
        title = slide_def.get("title", "")

        if slide_type == "title":
            ppt.add_title_slide(title, slide_def.get("subtitle", ""))

        elif slide_type == "logo":
            ppt.add_logo_slide()

        elif slide_type == "image":
            ppt.add_image_slide(
                title,
                slide_def.get("image", ""),
                action_title=slide_def.get("action_title", ""),
                text=slide_def.get("text", slide_def.get("body", "")),
                layout=slide_def.get("layout", "auto"),
            )

        elif slide_type == "table":
            ppt.add_table_slide(
                title,
                slide_def.get("data"),
                action_title=slide_def.get("action_title", ""),
                columns=slide_def.get("columns"),
                max_rows=slide_def.get("max_rows", 15),
            )

        else:  # content
            builder = ppt.add_content_slide(title, slide_def.get("action_title", ""))
            if "body" in slide_def or "text" in slide_def:
                builder.body(slide_def.get("body", slide_def.get("text", "")))
            if "image" in slide_def:
                builder.image(slide_def["image"], slide_def.get("layout", "auto"))

    return ppt.save(output)


# =============================================================================
# MODULE EXPORTS
# =============================================================================

__all__ = [
    "PPTGenerator",
    "PPTConfig",
    "SlideBuilder",
    "quick_ppt",
    "detect_image_layout",
    "BRAND",
    "IMAGE_LAYOUTS",
]
